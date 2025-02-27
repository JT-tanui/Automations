import os
import tkinter as tk
from tkinter import filedialog, messagebox, simpledialog
from PIL import Image
import pdf2image
import tempfile
from pptx import Presentation
from pptx.util import Inches

class PDFToPPTConverter:
    def __init__(self, root=None):
        self.pdf_files = []
        self.temp_dir = tempfile.mkdtemp()
        
        if root:
            self.root = root
        else:
            self.root = tk.Tk()
            self.root.title("PDF to PowerPoint Converter")
            self.root.geometry("600x500")
            self.setup_ui()
    
    def setup_ui(self):
        # Create frames
        self.top_frame = tk.Frame(self.root)
        self.top_frame.pack(fill=tk.X, padx=10, pady=10)
        
        self.list_frame = tk.Frame(self.root)
        self.list_frame.pack(fill=tk.BOTH, expand=True, padx=10, pady=10)
        
        self.button_frame = tk.Frame(self.root)
        self.button_frame.pack(fill=tk.X, padx=10, pady=10)
        
        # Top frame buttons
        self.select_folder_btn = tk.Button(self.top_frame, text="Select Default Folder", command=self.select_default_folder)
        self.select_folder_btn.pack(side=tk.LEFT, padx=5)
        
        self.select_files_btn = tk.Button(self.top_frame, text="Select PDF Files", command=self.select_pdf_files)
        self.select_files_btn.pack(side=tk.LEFT, padx=5)
        
        # PDF List
        self.list_label = tk.Label(self.list_frame, text="Selected PDFs:")
        self.list_label.pack(anchor=tk.W)
        
        self.listbox_frame = tk.Frame(self.list_frame)
        self.listbox_frame.pack(fill=tk.BOTH, expand=True)
        
        self.scrollbar = tk.Scrollbar(self.listbox_frame)
        self.scrollbar.pack(side=tk.RIGHT, fill=tk.Y)
        
        self.pdf_listbox = tk.Listbox(self.listbox_frame)
        self.pdf_listbox.pack(fill=tk.BOTH, expand=True)
        
        self.pdf_listbox.config(yscrollcommand=self.scrollbar.set)
        self.scrollbar.config(command=self.pdf_listbox.yview)
        
        # Button frame
        self.remove_btn = tk.Button(self.button_frame, text="Remove Selected", command=self.remove_selected)
        self.remove_btn.pack(side=tk.LEFT, padx=5)
        
        self.clear_btn = tk.Button(self.button_frame, text="Clear All", command=self.clear_all)
        self.clear_btn.pack(side=tk.LEFT, padx=5)
        
        self.generate_btn = tk.Button(self.button_frame, text="Generate PowerPoint", command=self.generate_ppt)
        self.generate_btn.pack(side=tk.RIGHT, padx=5)
    
    def select_default_folder(self):
        """Look for PDFs in the 'cert' folder in the root project directory."""
        cert_folder = os.path.join(os.getcwd(), "cert")
        
        if os.path.exists(cert_folder) and os.path.isdir(cert_folder):
            self.add_pdfs_from_folder(cert_folder)
        else:
            messagebox.showinfo("Info", "The 'cert' folder does not exist in the project root. Please create it or select files manually.")
            folder_path = filedialog.askdirectory(title="Select Folder Containing PDFs")
            if folder_path:
                self.add_pdfs_from_folder(folder_path)
    
    def add_pdfs_from_folder(self, folder_path):
        """Add all PDFs from the specified folder to the list."""
        new_files = [os.path.join(folder_path, f) for f in os.listdir(folder_path) 
                    if f.lower().endswith('.pdf')]
        
        if not new_files:
            messagebox.showinfo("Info", "No PDF files found in the selected folder.")
            return
        
        # Add unique files
        for file_path in new_files:
            if file_path not in self.pdf_files:
                self.pdf_files.append(file_path)
                self.pdf_listbox.insert(tk.END, os.path.basename(file_path))
    
    def select_pdf_files(self):
        """Open file dialog to select PDF files."""
        file_paths = filedialog.askopenfilenames(
            title="Select PDF Files",
            filetypes=[("PDF Files", "*.pdf")]
        )
        
        if not file_paths:
            return
        
        # Add unique files
        for file_path in file_paths:
            if file_path not in self.pdf_files:
                self.pdf_files.append(file_path)
                self.pdf_listbox.insert(tk.END, os.path.basename(file_path))
    
    def remove_selected(self):
        """Remove selected PDFs from the list."""
        selection = self.pdf_listbox.curselection()
        if not selection:
            messagebox.showinfo("Info", "Please select a PDF to remove.")
            return
        
        # Remove from bottom to top to maintain indices
        for index in sorted(selection, reverse=True):
            del self.pdf_files[index]
            self.pdf_listbox.delete(index)
    
    def clear_all(self):
        """Clear all PDFs from the list."""
        self.pdf_files = []
        self.pdf_listbox.delete(0, tk.END)
    
    def generate_ppt(self):
        """Generate PowerPoint from selected PDFs."""
        if not self.pdf_files:
            messagebox.showinfo("Info", "Please add at least one PDF file.")
            return
        
        # Ask for output filename
        output_file = filedialog.asksaveasfilename(
            title="Save PowerPoint As",
            defaultextension=".pptx",
            filetypes=[("PowerPoint Presentation", "*.pptx")]
        )
        
        if not output_file:
            return
        
        try:
            # Create presentation
            prs = Presentation()
            
            # Process each PDF
            for pdf_path in self.pdf_files:
                self.add_pdf_to_ppt(prs, pdf_path)
            
            # Save presentation
            prs.save(output_file)
            messagebox.showinfo("Success", f"PowerPoint created successfully!\nSaved to: {output_file}")
        
        except Exception as e:
            messagebox.showerror("Error", f"An error occurred: {str(e)}")
    
    def add_pdf_to_ppt(self, prs, pdf_path):
        """Convert a PDF page to an image and add it to the presentation."""
        try:
            # Convert PDF to images
            images = pdf2image.convert_from_path(pdf_path)
            
            for img in images:
                # Save image temporarily
                img_path = os.path.join(self.temp_dir, f"temp_img_{os.path.basename(pdf_path)}_{id(img)}.png")
                img.save(img_path, "PNG")
                
                # Add slide
                slide = prs.slides.add_slide(prs.slide_layouts[6])  # Blank layout
                
                # Calculate image dimensions to fit slide
                slide_width = prs.slide_width
                slide_height = prs.slide_height
                
                # Add image
                pic = slide.shapes.add_picture(img_path, 0, 0)
                
                # Scale image to fit slide
                scale_width = slide_width / pic.width
                scale_height = slide_height / pic.height
                scale = min(scale_width, scale_height)
                
                pic.width = int(pic.width * scale)
                pic.height = int(pic.height * scale)
                
                # Center image
                pic.left = int((slide_width - pic.width) / 2)
                pic.top = int((slide_height - pic.height) / 2)
                
                # Clean up temp file
                os.remove(img_path)
                
        except Exception as e:
            print(f"Error processing {pdf_path}: {str(e)}")
            raise
    
    def run(self):
        """Start the application."""
        self.root.mainloop()

if __name__ == "__main__":
    converter = PDFToPPTConverter()
    converter.run()